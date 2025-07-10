"""
ファイル差分比較アプリ用のテストファイル作成スクリプト
.doc, .docx, .ppt, .pptx, .xlsx, .pdf ファイルの各種差分パターンを生成します
"""

import os
import shutil
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter


def create_word_doc(file_path: str, content: str):
    """Word文書を作成（.doc/.docx 対応）"""
    try:
        doc = Document()
        doc.add_paragraph(content)
        doc.save(file_path)
        print(f"✅ Word文書作成: {file_path}")
    except Exception as e:
        print(f"❌ Word文書作成エラー {file_path}: {e}")


def create_excel_file(file_path: str, content: str):
    """Excel ファイルを作成"""
    try:
        wb = Workbook()
        ws = wb.active
        ws.title = "テストシート"

        # コンテンツを複数行に分割してセルに入力
        lines = content.split("\n") if "\n" in content else [content]
        for i, line in enumerate(lines, 1):
            ws[f"A{i}"] = line

        # サンプルデータも追加
        ws["B1"] = "項目"
        ws["B2"] = "値"
        ws["C1"] = "テスト"
        ws["C2"] = content[:20] + "..." if len(content) > 20 else content

        wb.save(file_path)
        print(f"✅ Excel作成: {file_path}")
    except Exception as e:
        print(f"❌ Excel作成エラー {file_path}: {e}")


def create_powerpoint(file_path: str, title: str, content: str):
    """PowerPoint プレゼンテーションを作成"""
    try:
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツのレイアウト
        slide = prs.slides.add_slide(slide_layout)

        # タイトルを設定
        title_shape = slide.shapes.title
        title_shape.text = title

        # コンテンツを設定
        content_shape = slide.placeholders[1]
        content_shape.text = content

        prs.save(file_path)
        print(f"✅ PowerPoint作成: {file_path}")
    except Exception as e:
        print(f"❌ PowerPoint作成エラー {file_path}: {e}")


def create_pdf_file(file_path: str, content: str):
    """PDF ファイルを作成"""
    try:
        # ファイル名の差分テスト用に、内容は同じにする必要があるパターンを判定
        is_name_diff_test = "ファイル名差分" in file_path
        is_no_diff_test = "差分なし" in file_path

        # 差分なしとファイル名差分の場合、同一のPDFファイルを作成する
        if is_name_diff_test or is_no_diff_test:
            create_identical_pdf(file_path)
            return

        c = canvas.Canvas(file_path, pagesize=letter)

        # 日本語対応のため、内容を英語に変換
        content_map = {
            "元の内容": "Original Content",
            "変更された内容": "Modified Content",
            "削除されるファイル": "File to be Deleted",
            "追加したファイル": "Added File",
            "旧内容": "Old Content",
            "新内容": "New Content",
        }

        # 日本語を英語に変換
        english_content = content_map.get(content, f"Content: {content}")

        # 標準フォントを使用
        c.setFont("Helvetica-Bold", 14)
        y_position = 750

        # タイトル
        c.drawString(100, y_position, "PDF Test File")
        y_position -= 30

        # ファイル名情報
        file_name = file_path.split("/")[-1].replace("\\", "/")
        file_name_map = {
            "差分なし": "no_diff",
            "内容差分": "content_diff",
            "ファイル名差分_変更前": "name_diff_before",
            "ファイル名差分_変更後": "name_diff_after",
            "削除": "deleted",
            "追加": "added",
            "内容・ファイル名差分_変更前": "both_diff_before",
            "内容・ファイル名差分_変更後": "both_diff_after",
        }

        english_file_name = file_name
        for japanese, english in file_name_map.items():
            english_file_name = english_file_name.replace(japanese, english)

        c.setFont("Helvetica", 12)
        c.drawString(100, y_position, f"File: {english_file_name}")
        y_position -= 25

        # コンテンツ
        c.setFont("Helvetica-Bold", 12)
        c.drawString(100, y_position, "Content:")
        y_position -= 20

        c.setFont("Helvetica", 12)
        lines = (
            english_content.split("\n")
            if "\n" in english_content
            else [english_content]
        )
        for line in lines:
            c.drawString(120, y_position, line)
            y_position -= 20

        # 作成日時
        y_position -= 10
        c.setFont("Helvetica", 10)
        from datetime import datetime

        c.drawString(
            100, y_position, f"Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        )

        # フッター
        c.drawString(100, 50, "Test file for file comparison application")

        c.save()
        print(f"✅ PDF作成: {file_path}")
    except Exception as e:
        print(f"❌ PDF作成エラー {file_path}: {e}")


def create_identical_pdf(file_path: str):
    """差分なし・ファイル名差分テスト用の同一PDFファイルを作成"""
    c = canvas.Canvas(file_path, pagesize=letter)

    # 完全に同一の内容でPDFを作成
    c.setFont("Helvetica-Bold", 14)
    c.drawString(100, 750, "PDF Test File")

    c.setFont("Helvetica", 12)
    c.drawString(100, 720, "File: identical_test.pdf")

    c.setFont("Helvetica-Bold", 12)
    c.drawString(100, 695, "Content:")

    c.setFont("Helvetica", 12)
    c.drawString(120, 675, "Same Content")

    c.setFont("Helvetica", 10)
    c.drawString(100, 640, "Created: 2025-01-01 00:00:00")

    c.drawString(100, 50, "Test file for file comparison application")

    c.save()
    print(f"✅ PDF作成: {file_path}")


def create_file_by_type(file_path: str, file_type: str, content: str):
    """ファイル形式に応じてファイルを作成"""
    try:
        if file_type == "word":
            create_word_doc(file_path, content)
        elif file_type == "excel":
            create_excel_file(file_path, content)
        elif file_type == "powerpoint":
            create_powerpoint(file_path, "テストファイル", content)
        elif file_type == "pdf":
            create_pdf_file(file_path, content)
        else:
            print(f"⚠️  未対応のファイル形式: {file_type}")
    except Exception as e:
        print(f"❌ ファイル作成エラー {file_path}: {e}")


def create_test_files():
    """テスト用のファイルを作成（ベースファイルを使わず、変更前→変更後にコピーして差分加工）"""
    print("🗂️  テストファイルを作成中...")

    # ディレクトリを作成
    for dir_name in ["変更前ディレクトリ", "変更後ディレクトリ", "保存先ディレクトリ"]:
        if not os.path.exists(dir_name):
            os.makedirs(dir_name)

    file_patterns = {
        "doc": {"ext": "doc", "type": "word"},
        "docx": {"ext": "docx", "type": "word"},
        "pptx": {"ext": "pptx", "type": "powerpoint"},
        "xlsx": {"ext": "xlsx", "type": "excel"},
        "pdf": {"ext": "pdf", "type": "pdf"},
    }

    # 1. 変更前ディレクトリに全パターンのファイルを作成
    for file_key, file_config in file_patterns.items():
        ext = file_config["ext"]
        file_type = file_config["type"]
        # 差分なし
        create_file_by_type(
            f"変更前ディレクトリ/{file_key}_差分なし.{ext}", file_type, "同じ内容"
        )
        # 内容差分
        create_file_by_type(
            f"変更前ディレクトリ/{file_key}_内容差分.{ext}", file_type, "元の内容"
        )
        # 削除（pdfのみ英語、それ以外は日本語）
        if file_type == "pdf":
            delete_content = "File to be Deleted - PDF Version"
        else:
            delete_content = "これは削除専用のダミーファイルです - 日本語版"
        create_file_by_type(
            f"変更前ディレクトリ/{file_key}_削除.{ext}", file_type, delete_content
        )
        # ファイル名差分_変更前
        create_file_by_type(
            f"変更前ディレクトリ/{file_key}_ファイル名差分_変更前.{ext}",
            file_type,
            "同じ内容",
        )
        # 内容・ファイル名差分_変更前
        if file_type == "pdf":
            both_diff_before_content = "Old Content (before rename and content change)"
        else:
            both_diff_before_content = "旧内容"
        create_file_by_type(
            f"変更前ディレクトリ/{file_key}_内容・ファイル名差分_変更前.{ext}",
            file_type,
            both_diff_before_content,
        )

    # 2. 変更後ディレクトリに必要なファイルをコピー
    for file_key, file_config in file_patterns.items():
        ext = file_config["ext"]
        file_type = file_config["type"]
        # 差分なし
        shutil.copy2(
            f"変更前ディレクトリ/{file_key}_差分なし.{ext}",
            f"変更後ディレクトリ/{file_key}_差分なし.{ext}",
        )
        # 内容差分（コピー後に内容変更）
        shutil.copy2(
            f"変更前ディレクトリ/{file_key}_内容差分.{ext}",
            f"変更後ディレクトリ/{file_key}_内容差分.{ext}",
        )
        create_file_by_type(
            f"変更後ディレクトリ/{file_key}_内容差分.{ext}", file_type, "変更された内容"
        )
        # ファイル名差分_変更前→ファイル名差分_変更後（名前変更のみ）
        shutil.copy2(
            f"変更前ディレクトリ/{file_key}_ファイル名差分_変更前.{ext}",
            f"変更後ディレクトリ/{file_key}_ファイル名差分_変更後.{ext}",
        )
        # 内容・ファイル名差分_変更前→内容・ファイル名差分_変更後（名前＋内容変更）
        shutil.copy2(
            f"変更前ディレクトリ/{file_key}_内容・ファイル名差分_変更前.{ext}",
            f"変更後ディレクトリ/{file_key}_内容・ファイル名差分_変更後.{ext}",
        )
        # PDFのみ内容を変更
        if file_type == "pdf":
            create_file_by_type(
                f"変更後ディレクトリ/{file_key}_内容・ファイル名差分_変更後.{ext}",
                file_type,
                "Changed Content (after rename and content change)",
            )
        else:
            create_file_by_type(
                f"変更後ディレクトリ/{file_key}_内容・ファイル名差分_変更後.{ext}",
                file_type,
                "新内容",
            )
    # 追加（変更後ディレクトリにのみ存在, 英語）
    for file_key, file_config in file_patterns.items():
        ext = file_config["ext"]
        file_type = file_config["type"]
        create_file_by_type(
            f"変更後ディレクトリ/{file_key}_追加.{ext}",
            file_type,
            "This is a dummy file for addition only - English Version",
        )

    print("\n✅ テストファイルの作成が完了しました！")
    print_test_summary()


def create_master_pdf_template(file_path: str):
    """PDF用のマスターテンプレートを作成（差分なし・ファイル名差分テスト用）"""
    c = canvas.Canvas(file_path, pagesize=letter)

    # 完全に同一の内容でPDFを作成
    c.setFont("Helvetica-Bold", 14)
    c.drawString(100, 750, "PDF Test File")

    c.setFont("Helvetica", 12)
    c.drawString(100, 720, "File: identical_test.pdf")

    c.setFont("Helvetica-Bold", 12)
    c.drawString(100, 695, "Content:")

    c.setFont("Helvetica", 12)
    c.drawString(120, 675, "Same Content")

    c.setFont("Helvetica", 10)
    c.drawString(100, 640, "Created: 2025-01-01 00:00:00")

    c.drawString(100, 50, "Test file for file comparison application")

    c.save()


def print_test_summary():
    """作成されたテストファイルのサマリーを表示"""
    print("\n📋 作成されたテストファイル:")
    print("=" * 60)

    for dir_name in ["変更前ディレクトリ", "変更後ディレクトリ"]:
        if os.path.exists(dir_name):
            files = os.listdir(dir_name)
            print(f"\n📂 {dir_name} ({len(files)}ファイル):")
            for file in sorted(files):
                print(f"  📄 {file}")

    print("\n🎯 期待される差分検出結果:")
    print("  ✅ 差分なし: 10ファイル (.doc, .docx, .pptx, .xlsx, .pdf)")
    print("  📝 内容変更: 10ファイル")
    print("  🔄 ファイル名変更: 10ファイル")
    print("  ➕ 追加: 5ファイル")
    print("  ➖ 削除: 5ファイル")
    print("  ⚠️  内容・ファイル名変更: 10ファイル")


if __name__ == "__main__":
    try:
        create_test_files()
        print("\n🚀 テストファイル作成完了！アプリを起動してテストしてください。")
    except Exception as e:
        print(f"❌ テストファイル作成中にエラーが発生しました: {e}")
        import traceback

        traceback.print_exc()
