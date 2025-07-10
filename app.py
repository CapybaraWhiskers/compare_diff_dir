import streamlit as st
import os
import shutil
from pathlib import Path
import hashlib
from typing import Dict, List, Tuple, Optional
import pandas as pd
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
from pdfminer.high_level import extract_text
import tempfile
import subprocess
import json

# ページ設定
st.set_page_config(page_title="ファイル差分比較アプリ", page_icon="📄", layout="wide")


class FileComparator:
    """ファイル比較クラス"""

    def __init__(self):
        self.supported_extensions = {".doc", ".docx", ".ppt", ".pptx", ".xlsx", ".pdf"}

    def get_files_in_directory(self, directory: str) -> Dict[str, str]:
        """ディレクトリ内のサポートされているファイル一覧を取得"""
        files = {}
        if not os.path.exists(directory):
            return files

        for file_path in Path(directory).rglob("*"):
            if (
                file_path.is_file()
                and file_path.suffix.lower() in self.supported_extensions
            ):
                relative_path = str(file_path.relative_to(directory))
                files[relative_path] = str(file_path)
        return files

    def extract_text_from_file(self, file_path: str) -> str:
        """ファイルからテキストを抽出"""
        try:
            file_ext = Path(file_path).suffix.lower()

            if file_ext in [".doc", ".docx"]:
                return self._extract_from_docx(file_path)
            elif file_ext in [".ppt", ".pptx"]:
                return self._extract_from_pptx(file_path)
            elif file_ext == ".xlsx":
                return self._extract_from_xlsx(file_path)
            elif file_ext == ".pdf":
                return self._extract_from_pdf(file_path)
            else:
                return ""
        except Exception as e:
            st.warning(
                f"ファイル {file_path} の読み込みでエラーが発生しました: {str(e)}"
            )
            return ""

    def _extract_from_docx(self, file_path: str) -> str:
        """DOCXファイルからテキストを抽出"""
        doc = Document(file_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return "\n".join(text)

    def _extract_from_pptx(self, file_path: str) -> str:
        """PPTXファイルからテキストを抽出"""
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _extract_from_xlsx(self, file_path: str) -> str:
        """XLSXファイルからテキストを抽出"""
        df = pd.read_excel(file_path, sheet_name=None)
        text = []
        for sheet_name, sheet_df in df.items():
            text.append(f"Sheet: {sheet_name}")
            text.append(sheet_df.to_string())
        return "\n".join(text)

    def _extract_from_pdf(self, file_path: str) -> str:
        """PDFファイルからテキストを抽出"""
        try:
            # PyMuPDFを使用
            doc = fitz.open(file_path)
            text = []
            for page in doc:
                text.append(page.get_text())
            doc.close()
            return "\n".join(text)
        except:
            # フォールバックとしてpdfminer.sixを使用
            return extract_text(file_path)

    def calculate_file_hash(self, file_path: str) -> str:
        """ファイルのハッシュ値を計算"""
        hash_md5 = hashlib.md5()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
        except Exception:
            return ""
        return hash_md5.hexdigest()

    def compare_directories(self, dir1: str, dir2: str) -> Dict:
        """ディレクトリ間のファイル差分を比較"""
        files1 = self.get_files_in_directory(dir1)
        files2 = self.get_files_in_directory(dir2)

        result = {
            "added": [],  # 追加されたファイル
            "deleted": [],  # 削除されたファイル
            "modified": [],  # 内容が変更されたファイル
            "renamed": [],  # ファイル名が変更されたファイル
            "unchanged": [],  # 変更なし
        }

        # 進捗バーの表示
        progress_bar = st.progress(0)
        status_text = st.empty()

        total_files = len(set(files1.keys()) | set(files2.keys()))
        current_file = 0

        # 追加されたファイル
        for file_name in files2.keys() - files1.keys():
            result["added"].append(
                {"name": file_name, "path": files2[file_name], "type": "added"}
            )
            current_file += 1
            progress_bar.progress(current_file / total_files)
            status_text.text(f"処理中... {current_file}/{total_files}")

        # 削除されたファイル
        for file_name in files1.keys() - files2.keys():
            result["deleted"].append(
                {"name": file_name, "path": files1[file_name], "type": "deleted"}
            )
            current_file += 1
            progress_bar.progress(current_file / total_files)
            status_text.text(f"処理中... {current_file}/{total_files}")

        # 共通ファイルの内容比較
        common_files = files1.keys() & files2.keys()
        for file_name in common_files:
            file1_path = files1[file_name]
            file2_path = files2[file_name]

            # ハッシュ値で高速比較
            hash1 = self.calculate_file_hash(file1_path)
            hash2 = self.calculate_file_hash(file2_path)

            if hash1 != hash2:
                # 差分詳細を取得
                try:
                    text1 = self.extract_text_from_file(file1_path)
                    text2 = self.extract_text_from_file(file2_path)
                    diff_lines = self._diff_text_lines(text1, text2)
                    text_similarity = self._calculate_text_similarity(text1, text2)
                    diff_summary = [
                        f"{n}行目:\n  変更前: {l1}\n  変更後: {l2}"
                        for n, l1, l2 in diff_lines
                    ]
                    # 類似度情報を追加（UI非表示化のためコメントアウト）
                    # if text_similarity < 0.3:  # 類似度30%未満
                    #     diff_summary.insert(
                    #         0,
                    #         f"⚠️ 内容の類似度が低いです（{text_similarity:.1%}）- 全く異なる内容の可能性があります",
                    #     )
                    # elif text_similarity < 0.7:  # 類似度70%未満
                    #     diff_summary.insert(
                    #         0, f"📝 内容の類似度: {text_similarity:.1%}"
                    #     )
                    # else:
                    #     diff_summary.insert(
                    #         0, f"✅ 内容の類似度: {text_similarity:.1%}"
                    #     )
                except Exception as e:
                    diff_lines = []
                    text_similarity = 0.0
                    diff_summary = [f"差分抽出エラー: {e}"]
                result["modified"].append(
                    {
                        "name": file_name,
                        "path1": file1_path,
                        "path2": file2_path,
                        "type": "modified",
                        "text_similarity": text_similarity,
                        "diff_lines": diff_lines,
                        "diff_summary": diff_summary,
                    }
                )
            else:
                result["unchanged"].append(
                    {
                        "name": file_name,
                        "path1": file1_path,
                        "path2": file2_path,
                        "type": "unchanged",
                    }
                )

            current_file += 1
            progress_bar.progress(current_file / total_files)
            status_text.text(f"処理中... {current_file}/{total_files}")

        # ファイル名変更の検出（内容ベース）
        try:
            self._detect_renamed_files(result, files1, files2)
        except Exception as e:
            st.warning(f"ファイル名変更検出でエラーが発生しました: {str(e)}")

        progress_bar.progress(1.0)
        status_text.text("比較完了!")

        return result

    def _diff_text_lines(self, text1: str, text2: str) -> List[Tuple[int, str, str]]:
        """2つのテキストの差分行を返す（行番号, 変更前, 変更後）"""
        lines1 = text1.splitlines()
        lines2 = text2.splitlines()
        max_len = max(len(lines1), len(lines2))
        diff = []
        for i in range(max_len):
            l1 = lines1[i] if i < len(lines1) else ""
            l2 = lines2[i] if i < len(lines2) else ""
            if l1 != l2:
                diff.append((i + 1, l1, l2))  # 1始まり
        return diff

    def _calculate_text_similarity(self, text1: str, text2: str) -> float:
        """2つのテキストの類似度を計算（0.0-1.0）"""
        if not text1 and not text2:
            return 1.0
        if not text1 or not text2:
            return 0.0

        # Levenshtein距離を使用した類似度計算
        def levenshtein_distance(s1: str, s2: str) -> int:
            if len(s1) < len(s2):
                return levenshtein_distance(s2, s1)
            if len(s2) == 0:
                return len(s1)

            previous_row = list(range(len(s2) + 1))
            for i, c1 in enumerate(s1):
                current_row = [i + 1]
                for j, c2 in enumerate(s2):
                    insertions = previous_row[j + 1] + 1
                    deletions = current_row[j] + 1
                    substitutions = previous_row[j] + (c1 != c2)
                    current_row.append(min(insertions, deletions, substitutions))
                previous_row = current_row

            return previous_row[-1]

        # 正規化（小文字・空白除去）
        normalized_text1 = "".join(text1.lower().split())
        normalized_text2 = "".join(text2.lower().split())

        max_len = max(len(normalized_text1), len(normalized_text2))
        if max_len == 0:
            return 1.0

        distance = levenshtein_distance(normalized_text1, normalized_text2)
        similarity = 1.0 - (distance / max_len)
        return max(0.0, similarity)

    def _detect_renamed_files(self, result: Dict, files1: Dict, files2: Dict):
        """ファイル名変更の検出（改善版）"""
        # 追加・削除されたファイルの中で内容が同じものを探す
        added_files = result["added"][:]  # コピーを作成
        deleted_files = result["deleted"][:]  # コピーを作成

        files_to_remove_from_added = []
        files_to_remove_from_deleted = []

        # ファイルサイズも考慮した高速マッチング
        added_file_info = []
        deleted_file_info = []

        # 追加ファイルの情報を収集
        for added_file in added_files:
            if added_file in files_to_remove_from_added:
                continue
            try:
                file_path = added_file["path"]
                file_size = os.path.getsize(file_path)
                file_hash = self.calculate_file_hash(file_path)
                if file_hash:
                    added_file_info.append(
                        {
                            "file": added_file,
                            "size": file_size,
                            "hash": file_hash,
                            "ext": Path(file_path).suffix.lower(),
                        }
                    )
            except Exception:
                continue

        # 削除ファイルの情報を収集
        for deleted_file in deleted_files:
            if deleted_file in files_to_remove_from_deleted:
                continue
            try:
                file_path = deleted_file["path"]
                file_size = os.path.getsize(file_path)
                file_hash = self.calculate_file_hash(file_path)
                if file_hash:
                    deleted_file_info.append(
                        {
                            "file": deleted_file,
                            "size": file_size,
                            "hash": file_hash,
                            "ext": Path(file_path).suffix.lower(),
                        }
                    )
            except Exception:
                continue

        # マッチング処理（ハッシュ + サイズ + 拡張子で判定）
        for deleted_info in deleted_file_info:
            if deleted_info["file"] in files_to_remove_from_deleted:
                continue

            for added_info in added_file_info:
                if added_info["file"] in files_to_remove_from_added:
                    continue

                # ハッシュ、サイズ、拡張子が全て一致する場合
                if (
                    deleted_info["hash"] == added_info["hash"]
                    and deleted_info["size"] == added_info["size"]
                    and deleted_info["ext"] == added_info["ext"]
                ):

                    # ファイル名変更として認識
                    result["renamed"].append(
                        {
                            "old_name": deleted_info["file"]["name"],
                            "new_name": added_info["file"]["name"],
                            "old_path": deleted_info["file"]["path"],
                            "new_path": added_info["file"]["path"],
                            "type": "renamed",
                        }
                    )

                    # 削除対象リストに追加
                    files_to_remove_from_added.append(added_info["file"])
                    files_to_remove_from_deleted.append(deleted_info["file"])
                    break

        # リストから削除（安全に削除）
        for file_item in files_to_remove_from_added:
            if file_item in result["added"]:
                result["added"].remove(file_item)

        for file_item in files_to_remove_from_deleted:
            if file_item in result["deleted"]:
                result["deleted"].remove(file_item)


class GitFileComparator(FileComparator):
    """Git機能を使ったファイル比較クラス（--no-indexオプション使用）"""

    def __init__(self):
        super().__init__()

    def is_git_available(self) -> bool:
        """Gitコマンドが利用可能かどうかを確認"""
        try:
            result = subprocess.run(
                ["git", "--version"],
                capture_output=True,
                text=True,
                timeout=10,
            )
            return result.returncode == 0
        except (subprocess.TimeoutExpired, FileNotFoundError):
            return False

    def detect_moved_and_modified_files_no_index(
        self, dir1: str, dir2: str, similarity_threshold: int = 50
    ) -> List[Dict]:
        """git diff --no-index を使用して内容と名前が同時に変更されたファイルを検出"""
        try:
            if not self.is_git_available():
                return []

            # git diff --no-index --find-renames を使用
            result = subprocess.run(
                [
                    "git",
                    "diff",
                    "--no-index",
                    f"--find-renames={similarity_threshold}",
                    "--name-status",
                    dir1,
                    dir2,
                ],
                capture_output=True,
                timeout=30,
                encoding=None,  # バイナリモードで取得
            )

            # git diff --no-index は差分がある場合は exitcode 1 を返すので、0 または 1 を正常とする
            if result.returncode not in [0, 1]:
                return []

            # エンコーディングを試行してデコード
            stdout = ""
            used_encoding = ""
            for encoding in ["utf-8", "cp932", "shift_jis", "latin1"]:
                try:
                    decoded = result.stdout.decode(encoding, errors="replace")
                    # 明らかに文字化けしている場合はスキップ
                    if decoded and decoded.count("\ufffd") < len(decoded) * 0.1:
                        stdout = decoded
                        used_encoding = encoding
                        break
                except (UnicodeDecodeError, AttributeError):
                    continue

            if not stdout:
                return []

            moved_and_modified = []
            debug_info = [f"使用エンコーディング: {used_encoding}"]

            # 警告行を除去して処理
            lines = [
                line
                for line in stdout.strip().split("\n")
                if line.strip() and not line.strip().startswith("warning:")
            ]

            debug_info.append(f"処理対象行数: {len(lines)}")

            for line_num, line in enumerate(lines, 1):
                if not line.strip():
                    continue

                debug_info.append(f"行{line_num}: {repr(line[:200])}")

                parts = line.split("\t")
                if len(parts) < 3:
                    debug_info.append(f"  → タブ分割不足: {len(parts)}個")
                    continue

                status = parts[0].strip()
                raw_old_path = parts[1].strip()
                raw_new_path = parts[2].strip()

                debug_info.append(f"  ステータス: {status}")
                debug_info.append(
                    f"  生パス: {repr(raw_old_path)} → {repr(raw_new_path)}"
                )

                # R で始まる場合（リネーム）
                if status.startswith("R"):
                    # 8進数エスケープシーケンスを処理
                    old_path = self._decode_git_path(raw_old_path)
                    new_path = self._decode_git_path(raw_new_path)

                    debug_info.append(f"  デコード後: {old_path} → {new_path}")

                    # パスの正規化
                    try:
                        old_path = os.path.normpath(old_path)
                        new_path = os.path.normpath(new_path)
                        debug_info.append(f"  正規化後: {old_path} → {new_path}")
                    except Exception as e:
                        debug_info.append(f"  正規化エラー: {e}")
                        continue

                    # ディレクトリプレフィックスを除去してファイル名を取得
                    old_name = self._extract_filename(old_path, dir1)
                    new_name = self._extract_filename(new_path, dir2)

                    debug_info.append(f"  抽出ファイル名: {old_name} → {new_name}")

                    # 不正なファイル名をスキップ
                    if (
                        not old_name
                        or not new_name
                        or len(old_name) > 255
                        or len(new_name) > 255
                        or old_name in [".", ".."]
                        or new_name in [".", ".."]
                    ):
                        debug_info.append(f"  → 不正ファイル名のためスキップ")
                        continue

                    # サポート対象のファイル形式のみを対象とする
                    old_ext = os.path.splitext(old_name)[1].lower()
                    new_ext = os.path.splitext(new_name)[1].lower()

                    debug_info.append(f"  拡張子: {old_ext} → {new_ext}")

                    if (
                        old_ext in self.supported_extensions
                        and new_ext in self.supported_extensions
                    ):

                        similarity = 100  # デフォルト値
                        if len(status) > 1 and status[1:].isdigit():
                            similarity = int(status[1:])

                        # 類似度が100%未満、または「内容・ファイル名差分」ファイルの場合
                        is_content_name_diff = (
                            "内容・ファイル名差分" in old_name
                            and "内容・ファイル名差分" in new_name
                        )

                        debug_info.append(
                            f"  類似度: {similarity}%, 内容差分ファイル: {is_content_name_diff}"
                        )

                        if similarity < 100 or is_content_name_diff:
                            debug_info.append(f"  → 検出対象として追加")
                            moved_and_modified.append(
                                {
                                    "old_name": old_name,
                                    "new_name": new_name,
                                    "old_path": old_path,
                                    "new_path": new_path,
                                    "similarity": similarity,
                                    "type": "renamed_and_modified",
                                    "debug_info": debug_info[
                                        -10:
                                    ],  # 最後の10行のデバッグ情報
                                }
                            )
                        else:
                            debug_info.append(f"  → 条件に合わないためスキップ")
                    else:
                        debug_info.append(f"  → サポート外拡張子のためスキップ")

            # デバッグ情報をセッション状態に保存
            if hasattr(st, "session_state"):
                st.session_state.git_debug_info = debug_info

            return moved_and_modified

        except (subprocess.TimeoutExpired, FileNotFoundError, Exception) as e:
            st.warning(f"Git差分検出でエラーが発生しました: {str(e)}")
            return []

    def _decode_git_path(self, path: str) -> str:
        """Git出力のパスをデコード（8進数エスケープシーケンス対応）"""
        # クォートを除去
        path = path.strip('"')

        # 8進数エスケープシーケンス（\345\244\211など）をデコード
        try:
            # バックスラッシュエスケープされた8進数を実際のバイトに変換
            import re

            def octal_to_byte(match):
                octal_str = match.group(1)
                return bytes([int(octal_str, 8)])

            # \nnn形式の8進数エスケープシーケンスを検出
            octal_pattern = r"\\(\d{3})"

            # エスケープシーケンスをバイト列に変換
            byte_parts = []
            last_end = 0

            for match in re.finditer(octal_pattern, path):
                # マッチする前の部分を追加
                if match.start() > last_end:
                    byte_parts.append(path[last_end : match.start()].encode("utf-8"))

                # 8進数をバイトに変換して追加
                octal_str = match.group(1)
                byte_parts.append(bytes([int(octal_str, 8)]))
                last_end = match.end()

            # 最後の部分を追加
            if last_end < len(path):
                byte_parts.append(path[last_end:].encode("utf-8"))

            # すべてのバイト部分を結合
            if byte_parts:
                combined_bytes = b"".join(byte_parts)
                # UTF-8でデコード
                return combined_bytes.decode("utf-8", errors="replace")
            else:
                return path

        except Exception:
            # エラーが発生した場合は元のパスを返す
            return path

    def _extract_filename(self, file_path: str, base_dir: str) -> str:
        """ファイルパスからファイル名を抽出"""
        try:
            # 絶対パスに変換
            abs_file_path = (
                os.path.abspath(file_path)
                if not os.path.isabs(file_path)
                else file_path
            )
            abs_base_dir = os.path.abspath(base_dir)

            # ベースディレクトリからの相対パスを取得
            if abs_file_path.startswith(abs_base_dir):
                relative_path = os.path.relpath(abs_file_path, abs_base_dir)
                return relative_path
            else:
                # ベースディレクトリに含まれない場合はファイル名のみ
                return os.path.basename(file_path)
        except Exception:
            # エラーが発生した場合はファイル名のみ
            return os.path.basename(file_path)

    def compare_directories_with_git_no_index(
        self, dir1: str, dir2: str, similarity_threshold: int = 50
    ) -> Dict:
        """git diff --no-index を使用したディレクトリ比較"""
        # まず通常の比較を実行
        result = self.compare_directories(dir1, dir2)

        # Git機能が利用可能な場合、追加の検出を行う
        if self.is_git_available():
            try:
                # 内容と名前が同時に変更されたファイルを検出
                moved_and_modified = self.detect_moved_and_modified_files_no_index(
                    dir1, dir2, similarity_threshold
                )

                # 差分詳細を付与
                for item in moved_and_modified:
                    try:
                        text1 = self.extract_text_from_file(item["old_path"])
                        text2 = self.extract_text_from_file(item["new_path"])
                        diff_lines = self._diff_text_lines(text1, text2)
                        text_similarity = self._calculate_text_similarity(text1, text2)
                        diff_summary = [
                            f"{n}行目:\n  変更前: {l1}\n  変更後: {l2}"
                            for n, l1, l2 in diff_lines
                        ]
                        # 類似度情報を追加
                        # if text_similarity < 0.3:  # 類似度30%未満
                        #     diff_summary.insert(
                        #         0,
                        #         f"⚠️ 内容の類似度が低いです（{text_similarity:.1%}）- 全く異なる内容の可能性があります",
                        #     )
                        # elif text_similarity < 0.7:  # 類似度70%未満
                        #     diff_summary.insert(
                        #         0, f"📝 内容の類似度: {text_similarity:.1%}"
                        #     )
                        # else:
                        #     diff_summary.insert(
                        #         0, f"✅ 内容の類似度: {text_similarity:.1%}"
                        #     )
                        item["diff_lines"] = diff_lines
                        item["diff_summary"] = diff_summary
                        item["text_similarity"] = text_similarity
                    except Exception as e:
                        item["diff_lines"] = []
                        item["diff_summary"] = [f"差分抽出エラー: {e}"]
                        item["text_similarity"] = 0.0

                # 検出されたファイルを既存の結果から除去し、新しいカテゴリに追加
                if moved_and_modified:
                    if "renamed_and_modified" not in result:
                        result["renamed_and_modified"] = []

                    # 検出されたファイルを削除・追加リストから除去
                    for item in moved_and_modified:
                        # 削除リストから旧ファイルを除去
                        result["deleted"] = [
                            f
                            for f in result["deleted"]
                            if f["name"] != item["old_name"]
                        ]
                        # 追加リストから新ファイルを除去
                        result["added"] = [
                            f for f in result["added"] if f["name"] != item["new_name"]
                        ]

                    result["renamed_and_modified"].extend(moved_and_modified)

                # Git情報をメタデータとして追加
                result["git_info"] = {
                    "is_git_available": True,
                    "moved_and_modified_count": len(moved_and_modified),
                    "similarity_threshold": similarity_threshold,
                    "method": "no-index",
                }

            except Exception as e:
                st.warning(f"Git情報の取得中にエラーが発生しました: {str(e)}")
                result["git_info"] = {"is_git_available": False, "error": str(e)}
        else:
            result["git_info"] = {"is_git_available": False, "reason": "Git not found"}

        return result


def main():
    st.title("📄 ファイル差分比較アプリ")
    st.markdown("---")

    # ディレクトリパスの設定
    col1, col2 = st.columns(2)

    with col1:
        dir1 = st.text_input("変更前ディレクトリ", value="./変更前ディレクトリ")

    with col2:
        dir2 = st.text_input("変更後ディレクトリ", value="./変更後ディレクトリ")

    save_dir = st.text_input(
        "保存先ディレクトリ",
        value="./保存先ディレクトリ",
        help="ファイルをコピーする保存先ディレクトリのパスを入力してください。\n例:\n- ./出力フォルダ (相対パス)\n- C:\\Users\\username\\Documents\\output (絶対パス)\n- Z:\\共有フォルダ\\バックアップ (ネットワークパス)",
    )

    # 保存先ディレクトリの存在確認と作成確認
    if save_dir:
        try:
            abs_path = os.path.abspath(save_dir)
            if os.path.exists(save_dir):
                st.info(f"📁 保存先: `{abs_path}` （既存ディレクトリ）")
            else:
                st.warning(f"📁 保存先: `{abs_path}` （新規作成されます）")
        except Exception as e:
            st.error(f"⚠️ 無効なパス: {str(e)}")

    # ディレクトリの存在確認
    dir1_exists = os.path.exists(dir1)
    dir2_exists = os.path.exists(dir2)

    col1, col2 = st.columns(2)
    with col1:
        if dir1_exists:
            st.success(f"✅ {dir1} が見つかりました")
        else:
            st.error(f"❌ {dir1} が見つかりません")

    with col2:
        if dir2_exists:
            st.success(f"✅ {dir2} が見つかりました")
        else:
            st.error(f"❌ {dir2} が見つかりません")

    if not (dir1_exists and dir2_exists):
        st.warning("比較を開始するには、両方のディレクトリが存在する必要があります。")
        return

    # Git検出オプション
    use_git = st.checkbox(
        "🔧 Git機能を使用してファイル移動＋内容変更を検出",
        help="Gitコマンドを使用して、ファイル名変更と内容変更が同時に行われたファイルを検出できます（リポジトリ不要）",
    )

    # デバッグオプション
    debug_mode = st.checkbox(
        "🐛 デバッグモードを有効にする",
        help="Git差分検出の詳細情報を表示します（問題の診断用）",
    )

    # Git設定オプション
    similarity_threshold = 50
    if use_git:
        st.markdown("#### Git設定")
        col1, col2 = st.columns(2)

        with col1:
            similarity_threshold = st.slider(
                "リネーム検出の類似度閾値",
                min_value=1,
                max_value=100,
                value=50,
                step=1,
                help="この値以上の類似度があるファイルをリネームとして検出します。値が低いほど、より大きな変更があってもリネームとして検出されます。",
            )

        with col2:
            st.markdown(f"**現在の設定:** {similarity_threshold}%")
            if similarity_threshold <= 30:
                st.warning(
                    "⚠️ 類似度が低いと、関係のないファイルも誤検出する可能性があります"
                )
            elif similarity_threshold >= 90:
                st.info(
                    "📌 類似度が高いと、わずかな変更でもリネームとして検出されません"
                )
            else:
                st.success("✅ 推奨範囲の設定です")

    # 比較実行ボタン
    if st.button("🔍 ファイル差分比較を実行", type="primary"):
        with st.spinner("ファイルを比較中..."):
            if use_git:
                comparator = GitFileComparator()
                result = comparator.compare_directories_with_git_no_index(
                    dir1, dir2, similarity_threshold
                )
            else:
                comparator = FileComparator()
                result = comparator.compare_directories(dir1, dir2)

            # セッション状態に結果を保存
            st.session_state.comparison_result = result
            st.session_state.comparator = comparator
            st.session_state.debug_mode = debug_mode

    # 結果の表示
    if "comparison_result" in st.session_state:
        result = st.session_state.comparison_result
        comparator = st.session_state.comparator

        st.markdown("---")
        st.header("📊 比較結果")

        # Git情報の表示
        if "git_info" in result and result["git_info"]["is_git_available"]:
            git_info = result["git_info"]
            threshold = git_info.get("similarity_threshold", 50)
            method = git_info.get("method", "unknown")
            st.info(
                f"🔧 Git機能が有効です。ファイル移動+内容変更を検出しました。（類似度閾値: {threshold}%, 方式: {method}）"
            )
            if "moved_and_modified_count" in git_info:
                moved_count = git_info["moved_and_modified_count"]
                if moved_count > 0:
                    st.success(f"🔄📝 名前＋内容変更: {moved_count}ファイル検出")
                else:
                    st.info(
                        f"📋 設定された類似度閾値（{threshold}%）では、名前＋内容変更ファイルは検出されませんでした。"
                    )
        elif "git_info" in result and not result["git_info"]["is_git_available"]:
            reason = result["git_info"].get("reason", "不明")
            st.warning(f"⚠️ Git機能が利用できません: {reason}")

        # デバッグ情報の表示
        debug_mode = st.session_state.get("debug_mode", False)
        if debug_mode and "git_info" in result:
            with st.expander("🐛 Git機能デバッグ情報", expanded=True):
                st.json(result["git_info"])

                # 詳細なデバッグ情報を表示
                if hasattr(st.session_state, "git_debug_info"):
                    st.markdown("**Git差分処理の詳細ログ:**")
                    debug_text = "\n".join(
                        st.session_state.git_debug_info[-50:]
                    )  # 最後の50行
                    st.text_area("処理ログ", debug_text, height=300)

                # 実際の検出結果も表示
                if "renamed_and_modified" in result and result["renamed_and_modified"]:
                    st.markdown("**検出された名前＋内容変更ファイル:**")
                    for i, item in enumerate(result["renamed_and_modified"], 1):
                        st.text(
                            f"{i}. {item['old_name']} → {item['new_name']} (類似度: {item.get('similarity', 'N/A')}%)"
                        )
                        st.text(f"   旧パス: {item.get('old_path', 'N/A')}")
                        st.text(f"   新パス: {item.get('new_path', 'N/A')}")

                        # 個別デバッグ情報がある場合は表示
                        if "debug_info" in item:
                            with st.expander(f"ファイル {i} の処理詳細"):
                                st.text("\n".join(item["debug_info"]))
                        st.text("")
                else:
                    st.warning("名前＋内容変更ファイルが検出されませんでした。")

                    # デバッグ情報から原因を分析
                    if hasattr(st.session_state, "git_debug_info"):
                        st.markdown("**考えられる原因:**")
                        debug_info = st.session_state.git_debug_info

                        # 処理対象行数をチェック
                        processing_lines = [
                            line for line in debug_info if "処理対象行数:" in line
                        ]
                        if processing_lines:
                            st.text(f"• {processing_lines[-1]}")

                        # スキップされた理由を分析
                        skip_reasons = [
                            line for line in debug_info if "スキップ" in line
                        ]
                        if skip_reasons:
                            st.text("• スキップされた理由:")
                            for reason in skip_reasons[-5:]:  # 最後の5個
                                st.text(f"  - {reason}")

                        # エンコーディング情報
                        encoding_info = [
                            line for line in debug_info if "エンコーディング:" in line
                        ]
                        if encoding_info:
                            st.text(f"• {encoding_info[-1]}")

        # サマリー表示の列数を調整
        if "renamed_and_modified" in result and result["renamed_and_modified"]:
            col1, col2, col3, col4, col5, col6 = st.columns(6)
        else:
            col1, col2, col3, col4, col5 = st.columns(5)

        with col1:
            st.metric(
                "追加",
                len(result["added"]),
                delta=len(result["added"]) if result["added"] else None,
            )
        with col2:
            st.metric(
                "削除",
                len(result["deleted"]),
                delta=-len(result["deleted"]) if result["deleted"] else None,
            )
        with col3:
            st.metric("内容変更", len(result["modified"]))
        with col4:
            st.metric("名前変更", len(result["renamed"]))
        with col5:
            st.metric("変更なし", len(result["unchanged"]))

        # 名前＋内容変更の列を追加
        if "renamed_and_modified" in result and result["renamed_and_modified"]:
            with col6:
                st.metric("名前＋内容変更", len(result["renamed_and_modified"]))

        # フィルタリングオプション
        st.markdown("---")
        st.subheader("🔧 フィルタリング")

        # 利用可能なオプションを動的に構築
        available_options = ["追加", "削除", "内容変更", "名前変更", "変更なし"]
        default_options = ["追加", "削除", "内容変更", "名前変更"]

        if "renamed_and_modified" in result and result["renamed_and_modified"]:
            available_options.insert(-1, "名前＋内容変更")
            default_options.append("名前＋内容変更")

        filter_options = st.multiselect(
            "表示する差分の種類を選択:",
            available_options,
            default=default_options,
        )

        # 結果の詳細表示とファイル選択
        selected_files = []

        if "追加" in filter_options and result["added"]:
            st.markdown("### ➕ 追加されたファイル")
            for item in result["added"]:
                if st.checkbox(f"📄 {item['name']}", key=f"added_{item['name']}"):
                    selected_files.append(("added", item))

        if "削除" in filter_options and result["deleted"]:
            st.markdown("### ➖ 削除されたファイル")
            for item in result["deleted"]:
                st.write(f"🗑️ {item['name']}")

        if "内容変更" in filter_options and result["modified"]:
            st.markdown("### 📝 内容変更されたファイル")
            for item in result["modified"]:
                if st.checkbox(f"📄 {item['name']}", key=f"modified_{item['name']}"):
                    selected_files.append(("modified", item))
                # 差分行情報を表示
                if "diff_summary" in item and item["diff_summary"]:
                    with st.expander(f"差分詳細: {item['name']}"):
                        for diff in item["diff_summary"]:
                            st.markdown(
                                diff.replace("\n", "<br>"), unsafe_allow_html=True
                            )

        if "名前変更" in filter_options and result["renamed"]:
            st.markdown("### 🔄 名前変更されたファイル")
            for item in result["renamed"]:
                if st.checkbox(
                    f"📄 {item['old_name']} → {item['new_name']}",
                    key=f"renamed_{item['new_name']}",
                ):
                    selected_files.append(("renamed", item))

        # 名前＋内容変更されたファイルの表示
        if (
            "名前＋内容変更" in filter_options
            and "renamed_and_modified" in result
            and result["renamed_and_modified"]
        ):
            st.markdown("### 🔄📝 名前＋内容変更されたファイル")
            for item in result["renamed_and_modified"]:
                similarity_text = (
                    f" (類似度: {item['similarity']}%)" if "similarity" in item else ""
                )
                if st.checkbox(
                    f"📄 {item['old_name']} → {item['new_name']}{similarity_text}",
                    key=f"renamed_modified_{item['new_name']}",
                ):
                    selected_files.append(("renamed_and_modified", item))
                # 差分行情報を表示
                if "diff_summary" in item and item["diff_summary"]:
                    with st.expander(
                        f"差分詳細: {item['old_name']} → {item['new_name']}"
                    ):
                        for diff in item["diff_summary"]:
                            st.markdown(
                                diff.replace("\n", "<br>"), unsafe_allow_html=True
                            )

        if "変更なし" in filter_options and result["unchanged"]:
            st.markdown("### ✅ 変更なしのファイル")
            for item in result["unchanged"]:
                if st.checkbox(f"📄 {item['name']}", key=f"unchanged_{item['name']}"):
                    selected_files.append(("unchanged", item))

        # ファイルコピー機能
        if selected_files:
            st.markdown("---")
            st.subheader("💾 選択したファイルをコピー")

            # 保存先ディレクトリの確認
            if save_dir.strip():
                save_path_display = os.path.abspath(save_dir) if save_dir else ""
                st.info(f"📁 コピー先: `{save_path_display}`")

                if st.button("📁 選択したファイルを保存先にコピー", type="secondary"):
                    copy_files(selected_files, save_dir, dir2)
            else:
                st.warning("⚠️ 保存先ディレクトリを指定してください。")


def copy_files(selected_files: List[Tuple], save_dir: str, source_dir: str):
    """選択されたファイルを保存先にコピー"""
    try:
        # 保存先ディレクトリの絶対パスを取得
        abs_save_dir = os.path.abspath(save_dir)

        # 保存先ディレクトリを作成
        if not os.path.exists(abs_save_dir):
            os.makedirs(abs_save_dir)
            st.info(f"📁 保存先ディレクトリを作成しました: {abs_save_dir}")

        success_count = 0
        error_count = 0

        progress_bar = st.progress(0)
        status_text = st.empty()

        for i, (file_type, item) in enumerate(selected_files):
            try:
                if file_type == "added":
                    source_path = item["path"]
                    dest_path = os.path.join(abs_save_dir, item["name"])
                elif file_type == "modified":
                    source_path = item["path2"]  # 内容変更後のファイル
                    dest_path = os.path.join(abs_save_dir, item["name"])
                elif file_type == "renamed":
                    source_path = item["new_path"]
                    dest_path = os.path.join(abs_save_dir, item["new_name"])
                elif file_type == "renamed_and_modified":
                    # Git検出された名前+内容変更ファイルの場合
                    source_path = item["new_path"]  # 新しいファイルパス
                    dest_path = os.path.join(abs_save_dir, item["new_name"])
                elif file_type == "unchanged":
                    source_path = item["path2"]
                    dest_path = os.path.join(abs_save_dir, item["name"])

                # ディレクトリ構造を作成
                os.makedirs(os.path.dirname(dest_path), exist_ok=True)

                # ファイルをコピー
                shutil.copy2(source_path, dest_path)
                success_count += 1

            except Exception as e:
                st.error(
                    f"ファイルのコピーに失敗しました: {item.get('name', 'Unknown')} - {str(e)}"
                )
                error_count += 1

            progress_bar.progress((i + 1) / len(selected_files))
            status_text.text(f"コピー中... {i + 1}/{len(selected_files)}")

        progress_bar.progress(1.0)
        status_text.text("コピー完了!")

        if success_count > 0:
            st.success(f"✅ {success_count} 個のファイルが正常にコピーされました!")
            st.info(f"📁 保存先: {abs_save_dir}")

        if error_count > 0:
            st.error(f"❌ {error_count} 個のファイルでエラーが発生しました。")

    except Exception as e:
        st.error(f"保存先ディレクトリの作成に失敗しました: {str(e)}")
        st.info("💡 ヒント: 書き込み権限があるディレクトリを指定してください。")


if __name__ == "__main__":
    main()
